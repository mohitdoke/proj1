"""Render a single-company summary slide, cloned from the Fund 1 Newsletter
deck's per-company slide design (navy/gold IIFL theme).

The slide carries the three sections the dashboard shows, in the newsletter's
own layout:

    top-left      Business Description
    top-right     Key Performance Indicators   (last 2 FYs + last 2 quarters)
    bottom-right  Key Financial Highlights     (last 2 FYs + last 2 quarters)
    bottom-left   Performance Summary / Commentary — the page's narrative
                  bullets, its stat tiles, and any footnotes

Usage:
    from _generate_company_report import generate_report
    generate_report(content_dict, "Output.pptx")

`content_dict` is exactly what src/lib/reportContent.py's JS counterpart
(buildReportContent) produces — every figure in it is already a formatted
string, computed by the website's own calculation engine. NOTHING in this
module does arithmetic or decides which row means "revenue": it lays out text
that has already been decided elsewhere. Keeping that boundary is what stops
the deck and the dashboard from disagreeing.

    {
      "legal_entity_name": str, "brand": str,
      "period_labels": [str, ...],                 # 4 columns
      "business": {"title", "description", "tags": [...], "scaleMetrics": [...]},
      "kpi_section":       {"title", "cornerLabel", "rows": [{"label", "values": [...]}], "notes": [...]},
      "financial_section": {"title", "cornerLabel", "rows": [...], "notes": [...]},
      "summary": {"title", "eyebrow", "bullets": [str, ...], "tiles": [{"label","value","sub"}], "notes": [...]},
      "footnotes": [str, ...],
    }
"""
import copy
import os

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

SOURCE_DECK = os.path.join(os.path.dirname(__file__), "..", "..", "assets", "Fund 1 Newsletter - Q4FY26 (1).pptx")
TEMPLATE_SLIDE_INDEX = 2  # Leegality slide - the per-company template to clone

# shape.name -> role, stable across clones since names live in the copied XML
SHAPE_TITLE = "Title 1"
SHAPE_DESC_HEADING = "TextBox 2"
SHAPE_COMMENTARY_HEADING = "TextBox 3"
SHAPE_FIN_HEADING = "TextBox 4"
SHAPE_KPI_HEADING = "TextBox 5"
SHAPE_DESC_BODY = "Rectangle 6"
SHAPE_LOGO = "Picture 7"
SHAPE_KPI_TABLE = "Table 12"
SHAPE_FIN_TABLE = "Table 13"
SHAPE_COMMENTARY_BODY = "TextBox 10"
SHAPE_TITLE_RULE = "Straight Connector 20"  # the rule between the title and the logo

# The template's own body type sizes, and the floor auto-fit will shrink to
# rather than let a long narrative spill off the slide.
BODY_PT = 9.0
BODY_MIN_PT = 6.5
TABLE_PT = 10.5
TABLE_MIN_PT = 7.5


def _clone_slide(prs, source_slide):
    """Deep-copy a slide's shape tree onto a freshly added slide of the same layout,
    re-wiring image relationships so blip references resolve in the new slide."""
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # drop placeholders python-pptx auto-added from the layout - we bring our own copies
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # map old r:id -> new r:id for image/media relationships
    rid_map = {}
    for rel_id, rel in source_slide.part.rels.items():
        if "image" in rel.reltype:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel_id] = new_rid

    blip_tag = qn("a:blip")
    embed_attr = qn("r:embed")
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        for blip in new_el.iter(blip_tag):
            old_rid = blip.get(embed_attr)
            if old_rid in rid_map:
                blip.set(embed_attr, rid_map[old_rid])
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def _remove_all_other_slides(prs, keep_slide):
    """Strip every slide except `keep_slide` from the presentation's slide list."""
    xml_slides = prs.slides._sldIdLst
    keep_rid = next(rel_id for rel_id, rel in prs.part.rels.items() if rel.target_part is keep_slide.part)
    for sld in list(xml_slides):
        if sld.get(qn("r:id")) != keep_rid:
            xml_slides.remove(sld)
            prs.part.drop_rel(sld.get(qn("r:id")))


# ------------------------------------------------------------------
# TEXT — every write below reuses the run properties (font, size, colour,
# theme reference) already sitting in the template's own XML, instead of
# letting python-pptx create a bare run that falls back to the default
# 18pt black. That's what keeps a generated slide looking like the deck it
# was cloned from.
# ------------------------------------------------------------------
def _first_run_props(text_frame):
    """A deepcopy-able <a:rPr> from the first run that has one, or None."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn("a:rPr"))
            if rPr is not None:
                return copy.deepcopy(rPr)
    return None


def _clear_paragraphs(text_frame):
    """Leave exactly one empty paragraph, keeping its paragraph properties."""
    paragraphs = text_frame.paragraphs
    for para in paragraphs[1:]:
        para._p.getparent().remove(para._p)
    first = paragraphs[0]
    for child in list(first._p):
        if child.tag in (qn("a:r"), qn("a:br"), qn("a:fld")):
            first._p.remove(child)
    return first


def _add_run(paragraph, text, rPr_template, size_pt=None, bold=None):
    run = paragraph.add_run()
    run.text = text
    if rPr_template is not None:
        new_rPr = copy.deepcopy(rPr_template)
        # add_run() already created an rPr; replace it wholesale so the
        # template's font/colour/theme attributes carry over.
        old = run._r.find(qn("a:rPr"))
        if old is not None:
            run._r.remove(old)
        run._r.insert(0, new_rPr)
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    return run


_BULLET_TAGS = ("a:buChar", "a:buAutoNum", "a:buNone", "a:buFont", "a:buBlip", "a:buClr", "a:buSzPct", "a:buSzPts")
# Children that must follow the bullet properties in a valid <a:pPr>.
_AFTER_BULLET_TAGS = ("a:tabLst", "a:defRPr", "a:extLst")


def _flatten_paragraph(paragraph):
    """Drop any inherited bullet glyph and hanging indent.

    The template's body paragraphs carry a bullet char and a hanging indent
    (the source slide's description is a bulleted list). Content written here
    supplies its own "•" where it wants one, so an inherited glyph would show
    up doubled on some lines and not others, with the first line indented
    differently from every line added after it.
    """
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in _BULLET_TAGS:
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu_none = pPr.makeelement(qn("a:buNone"), {})
    anchor = next((pPr.find(qn(t)) for t in _AFTER_BULLET_TAGS if pPr.find(qn(t)) is not None), None)
    if anchor is not None:
        anchor.addprevious(bu_none)
    else:
        pPr.append(bu_none)


def _set_lines(shape, lines, size_pt=None):
    """Write `lines` into a text box, one paragraph each, in the template's own
    type. `lines` items are (text, bold) tuples or plain strings."""
    tf = shape.text_frame
    tf.word_wrap = True
    rPr = _first_run_props(tf)
    first = _clear_paragraphs(tf)

    para = first
    for i, line in enumerate(lines):
        text, bold = line if isinstance(line, tuple) else (line, None)
        if i > 0:
            para = tf.add_paragraph()
        _flatten_paragraph(para)
        _add_run(para, text, rPr, size_pt=size_pt, bold=bold)
    return tf


def _set_single_line(shape, text):
    """Replace a one-line label/title, keeping its existing run formatting."""
    _set_lines(shape, [text])


# ------------------------------------------------------------------
# AUTO-FIT — python-pptx can't measure text, so this estimates wrapped line
# count from the box width and shrinks the type until it fits. Deliberately
# conservative: it is better for a dense company's commentary to render at
# 7pt than to run off the bottom of the slide.
# ------------------------------------------------------------------
def _estimate_lines(text, chars_per_line):
    if not text:
        return 1
    words, lines, current = text.split(), 0, 0
    for word in words:
        add = len(word) + (1 if current else 0)
        if current + add > chars_per_line and current:
            lines += 1
            current = len(word)
        else:
            current += add
    return lines + 1


def _fit_size(lines, box_width_in, box_height_in, start_pt, min_pt, line_spacing=1.22):
    """Largest size in [min_pt, start_pt] whose estimated wrapped height fits."""
    size = start_pt
    while size > min_pt:
        # ~0.55 em average glyph advance for this deck's sans face.
        chars_per_line = max(12, int(box_width_in * 72 / (size * 0.55)))
        line_h_in = size * line_spacing / 72
        total = sum(_estimate_lines(t if isinstance(t, str) else t[0], chars_per_line) for t in lines)
        if total * line_h_in <= box_height_in:
            return size
        size -= 0.5
    return min_pt


EMU_PER_IN = 914400
# The template's own footer sits at 7.25" on a 7.5" slide, and its commentary
# box is drawn taller than the space actually left below it (harmless in the
# source deck, where the text is short). Text is fitted to the real usable
# height so a long narrative stops above the footer instead of running off.
TEXT_BOTTOM_IN = 7.15


def _shape_box_in(shape):
    """(width, usable height) in inches — height clamped to what's actually
    left above the footer, which can be less than the shape's own height."""
    top_in = shape.top / EMU_PER_IN
    height_in = min(shape.height / EMU_PER_IN, max(0.5, TEXT_BOTTOM_IN - top_in))
    return shape.width / EMU_PER_IN, height_in


# ------------------------------------------------------------------
# TABLES — the template's tables are fixed at 6 columns and a fixed row
# count. Both are reshaped here by copying the template's own <a:gridCol>
# and <a:tr> elements, so added rows/columns inherit the deck's borders,
# fills and banding instead of arriving unstyled.
# ------------------------------------------------------------------
def _set_column_count(table, n_cols, total_width):
    """Grow/shrink to n_cols, then redistribute width: a wide label column plus
    equal value columns."""
    grid = table._tbl.find(qn("a:tblGrid"))
    cols = grid.findall(qn("a:gridCol"))

    while len(cols) > n_cols:
        grid.remove(cols[-1])
        for tr in table._tbl.findall(qn("a:tr")):
            tcs = tr.findall(qn("a:tc"))
            tr.remove(tcs[-1])
        cols = grid.findall(qn("a:gridCol"))

    while len(cols) < n_cols:
        grid.append(copy.deepcopy(cols[-1]))
        for tr in table._tbl.findall(qn("a:tr")):
            tcs = tr.findall(qn("a:tc"))
            tr.append(copy.deepcopy(tcs[-1]))
        cols = grid.findall(qn("a:gridCol"))

    label_w = int(total_width * 0.34)
    value_w = int((total_width - label_w) / max(1, n_cols - 1))
    table.columns[0].width = label_w
    for j in range(1, n_cols):
        table.columns[j].width = value_w


def _set_row_count(table, n_rows):
    """Grow/shrink to n_rows (header included). New rows clone the FIRST body
    row, so they arrive with the template's own cell formatting and its
    standard row height — the template's last row is a taller one, and cloning
    that would stretch the table past its box."""
    tbl = table._tbl
    rows = tbl.findall(qn("a:tr"))
    template_row = rows[1] if len(rows) > 1 else rows[0]
    body_height = template_row.get("h")
    while len(rows) > n_rows:
        tbl.remove(rows[-1])
        rows = tbl.findall(qn("a:tr"))
    while len(rows) < n_rows:
        tbl.append(copy.deepcopy(template_row))
        rows = tbl.findall(qn("a:tr"))
    # Uniform body rows: the source deck's own last row is taller than the
    # rest, which reads as an accident once the row order changes.
    if body_height:
        for tr in rows[1:]:
            tr.set("h", body_height)


def _cell_run_props(table):
    """Run properties for a header cell and for a body cell, sampled from the
    template before anything is overwritten."""
    header = _first_run_props(table.cell(0, 0).text_frame)
    body = None
    if len(table.rows) > 1:
        for j in range(len(table.columns)):
            body = _first_run_props(table.cell(1, j).text_frame)
            if body is not None:
                break
    return header, body if body is not None else header


def _write_cell(cell, text, rPr, size_pt, bold):
    tf = cell.text_frame
    tf.word_wrap = True
    para = _clear_paragraphs(tf)
    _add_run(para, text, rPr, size_pt=size_pt, bold=bold)


def _fill_table(table_shape, corner_label, period_labels, data_rows):
    """Reshape the template table to exactly (1 + len(data_rows)) x
    (1 + len(period_labels)) and fill it.

    `data_rows` are dicts of {"label": str, "values": [str, ...]} — every value
    already formatted upstream. Label cells stay bold like the template's
    metric rows; growth/percentage sub-rows (the ones the page renders in a
    lighter weight under their metric) are detected by label and rendered
    unbolded, matching the deck.
    """
    table = table_shape.table
    header_rPr, body_rPr = _cell_run_props(table)

    n_cols = len(period_labels) + 1
    _set_column_count(table, n_cols, table_shape.width)
    _set_row_count(table, len(data_rows) + 1)

    # Shrink type if a company's labels are long enough to wrap badly.
    longest_label = max([len(corner_label)] + [len(r["label"]) for r in data_rows] or [0])
    label_col_in = table.columns[0].width / 914400
    size = TABLE_PT if longest_label * 0.055 <= label_col_in else max(TABLE_MIN_PT, TABLE_PT - 1.5)

    _write_cell(table.cell(0, 0), corner_label, header_rPr, size, True)
    for j, label in enumerate(period_labels, start=1):
        _write_cell(table.cell(0, j), str(label), header_rPr, size, True)

    for i, row in enumerate(data_rows, start=1):
        label = row["label"]
        # Growth and margin rows qualify the metric row above them rather than
        # being metrics in their own right — the source deck sets exactly those
        # lighter ("%Growth YoY" under "Net Revenue", "%Gross Margin" under
        # "Gross Margin"), so the same rule is applied here.
        low = label.lower().strip()
        is_subrow = "growth" in low or low.endswith("margin")
        _write_cell(table.cell(i, 0), label, body_rPr, size, not is_subrow)
        for j in range(1, n_cols):
            value = row["values"][j - 1] if j - 1 < len(row["values"]) else ""
            _write_cell(table.cell(i, j), str(value), body_rPr, size, not is_subrow)


"""Right edge the title may grow to when there is no company logo beside it —
short of the deck's own IIFL mark in the top-right corner."""
TITLE_RIGHT_LIMIT_IN = 10.0
TITLE_MIN_PT = 14.0


def _master_title_pt(slide, default=24.0):
    """The deck's own title size, from the master's title text style."""
    tx_styles = slide.slide_layout.slide_master.element.find(qn("p:txStyles"))
    if tx_styles is None:
        return default
    title_style = tx_styles.find(qn("p:titleStyle"))
    lvl1 = title_style.find(qn("a:lvl1pPr")) if title_style is not None else None
    def_rPr = lvl1.find(qn("a:defRPr")) if lvl1 is not None else None
    sz = def_rPr.get("sz") if def_rPr is not None else None
    return float(sz) / 100 if sz else default


def _fit_title(slide, title_shape, text):
    """Keep "{legal entity} - {brand}" on one line: some legal names are long
    enough to wrap onto a second line and cross the rule under the title."""
    base = _master_title_pt(slide)
    width_in = title_shape.width / EMU_PER_IN
    size = base
    while size > TITLE_MIN_PT and len(text) * size * 0.55 / 72 > width_in:
        size -= 1
    _set_lines(title_shape, [text], size_pt=None if size >= base else size)


def _apply_logo(slide, by_name, new_image_path, brand=""):
    """Swap in this company's logo, or clear the template's.

    The cloned slide carries the template company's logo and the rule that
    separates it from the title. Left alone, every generated deck would be
    stamped with another company's brand mark — so with no logo supplied both
    are removed and the title is given the freed width, unless the company is Leegality.
    """
    logo_shape = by_name.get(SHAPE_LOGO)
    if logo_shape is None:
        return
    left, top, width, height = logo_shape.left, logo_shape.top, logo_shape.width, logo_shape.height
    if new_image_path:
        logo_shape._element.getparent().remove(logo_shape._element)
        slide.shapes.add_picture(new_image_path, left, top, width, height)
        return
    if brand and "leegality" in brand.lower():
        # Keep the template's Leegality logo
        return
    logo_shape._element.getparent().remove(logo_shape._element)
    _remove_shape(by_name, SHAPE_TITLE_RULE)
    title = by_name.get(SHAPE_TITLE)
    if title is not None:
        title.width = max(title.width, int(TITLE_RIGHT_LIMIT_IN * EMU_PER_IN) - title.left)


def _remove_shape(by_name, name):
    shape = by_name.get(name)
    if shape is not None:
        shape._element.getparent().remove(shape._element)


# ------------------------------------------------------------------
# SECTION BUILDERS
# ------------------------------------------------------------------
def _business_lines(business):
    """Description paragraph, then the page's tag chips and scale metrics as
    compact single lines (there is no room for chip art at this size, and the
    words carry the same information)."""
    lines = []
    if business.get("description"):
        lines.append((business["description"], None))
    tags = business.get("tags") or []
    if tags:
        lines.append((" · ".join(tags), True))
    for metric in business.get("scaleMetrics") or []:
        lines.append((f"{metric['value']} — {metric['label']}", None))
    return lines or [("No business description is recorded for this company yet.", None)]


def _summary_lines(summary, footnotes):
    """Narrative bullets, then the page's stat tiles, then any footnote."""
    lines = []
    eyebrow = summary.get("eyebrow")
    if eyebrow:
        lines.append((eyebrow.upper(), True))
    for bullet in summary.get("bullets") or []:
        lines.append((f"• {bullet}", None))

    tiles = summary.get("tiles") or []
    if tiles:
        lines.append(("", None))
        for tile in tiles:
            sub = tile.get("sub")
            lines.append((f"{tile['label']}: {tile['value']}" + (f" ({sub})" if sub and sub != "—" else ""), True))

    for note in footnotes or []:
        lines.append((note, None))
    return lines


def generate_report(content: dict, output_path: str):
    prs = Presentation(SOURCE_DECK)
    template_slide = prs.slides[TEMPLATE_SLIDE_INDEX]
    new_slide = _clone_slide(prs, template_slide)
    by_name = {s.name: s for s in new_slide.shapes}

    brand = content.get("brand") or content.get("legal_entity_name") or "Company"
    legal = content.get("legal_entity_name") or brand
    if legal.lower().strip() == brand.lower().strip() or not legal:
        title_text = brand
    else:
        title_text = f"{legal} - {brand}"

    # Logo first: with none supplied it frees the width beside the title.
    _apply_logo(new_slide, by_name, content.get("logo_path"), brand)
    _fit_title(new_slide, by_name[SHAPE_TITLE], title_text)

    period_labels = content.get("period_labels") or []

    # --- Business Description (top-left) ---
    business = content.get("business") or {}
    _set_single_line(by_name[SHAPE_DESC_HEADING], business.get("title") or "Business Description")
    desc_shape = by_name[SHAPE_DESC_BODY]
    desc_lines = _business_lines(business)
    w, h = _shape_box_in(desc_shape)
    _set_lines(desc_shape, desc_lines, size_pt=_fit_size(desc_lines, w - 0.2, h - 0.15, BODY_PT, BODY_MIN_PT))

    # --- Key Performance Indicators (top-right) ---
    kpi = content.get("kpi_section")
    if kpi and kpi.get("rows"):
        _set_single_line(by_name[SHAPE_KPI_HEADING], kpi.get("title") or "Key Performance Indicators")
        _fill_table(by_name[SHAPE_KPI_TABLE], kpi.get("cornerLabel") or "KPI", period_labels, kpi["rows"])
    else:
        # No KPI table for this company's layout — drop the heading and the
        # table rather than leaving the template's Leegality figures on screen.
        _remove_shape(by_name, SHAPE_KPI_HEADING)
        _remove_shape(by_name, SHAPE_KPI_TABLE)

    # --- Key Financial Highlights (bottom-right) ---
    fin = content.get("financial_section")
    if fin and fin.get("rows"):
        _set_single_line(by_name[SHAPE_FIN_HEADING], fin.get("title") or "Key Financial Highlights")
        _fill_table(by_name[SHAPE_FIN_TABLE], fin.get("cornerLabel") or "INR Cr.", period_labels, fin["rows"])
    else:
        _remove_shape(by_name, SHAPE_FIN_HEADING)
        _remove_shape(by_name, SHAPE_FIN_TABLE)

    # --- Performance Summary / Commentary (bottom-left) ---
    summary = content.get("summary") or {}
    _set_single_line(by_name[SHAPE_COMMENTARY_HEADING], summary.get("title") or "Performance Summary")
    body_shape = by_name[SHAPE_COMMENTARY_BODY]
    footnotes = list(content.get("footnotes") or [])
    for section in (kpi, fin):
        if section:
            footnotes.extend(section.get("notes") or [])
    summary_lines = _summary_lines(summary, footnotes)
    w, h = _shape_box_in(body_shape)
    _set_lines(body_shape, summary_lines, size_pt=_fit_size(summary_lines, w - 0.2, h - 0.15, BODY_PT, BODY_MIN_PT))

    _remove_all_other_slides(prs, new_slide)

    prs.save(output_path)


if __name__ == "__main__":
    import json
    import sys
    if len(sys.argv) >= 3:
        with open(sys.argv[1], "r", encoding="utf-8") as f:
            data = json.load(f)
        generate_report(data, sys.argv[2])
    else:
        print("Usage: py _generate_company_report.py <input_json> <output_pptx>")
